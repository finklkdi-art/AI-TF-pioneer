from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import datetime

# --- 설정 및 데이터 ---
START_DATE = datetime.date(2026, 4, 21)
MILESTONE_DAYS = [2, 9, 16, 30]

# 폰트 설정 (시스템에 해당 폰트가 설치되어 있어야 정상 적용됩니다)
FONT_BOLD = 'Sandoll Gothic Neo1 B'
FONT_MEDIUM = 'Sandoll Gothic Neo1 M'
FONT_LIGHT = 'Sandoll Gothic Neo1 L'

milestones = [
    {
        "title": "[PHASE 0] AE Voice 웹 설문 조사",
        "date_offset": MILESTONE_DAYS[0],
        "mission": "전사 AE 대상 웹 설문을 통한 업무 프로세스 전수 조사 및 핵심 Pain Point 정량화.",
        "tasks": ["Google Forms/Typeform 활용 웹 설문 설계 및 배포", "견적 관련 예외 상황(Edge Case) 데이터 수집", "설문 결과 분석을 통한 기능 우선순위(Backlog) 도출"],
        "tech": ["Web Survey Tools", "Data Analysis (Excel/Pandas)"]
    },
    {
        "title": "[PHASE 1] 멀티 포맷 파서(Parser) 구축",
        "date_offset": MILESTONE_DAYS[1],
        "mission": "엑셀, PDF, 이미지(JPG), HWP 등 파편화된 외부 견적 데이터를 하나의 규격으로 통합.",
        "tasks": ["멀티 포맷(Excel/PDF/IMG/HWP) 대응 파싱 로직", "지능형 텍스트 추출(OCR) 엔진 연동", "외주비 데이터 표준화 및 정제"],
        "tech": ["Pandas", "PyMuPDF", "OCR Engine", "Python-docx"]
    },
    {
        "title": "[PHASE 2] 엑셀 기반 다이나믹 리포팅",
        "date_offset": MILESTONE_DAYS[2],
        "mission": "클라이언트별 맞춤형 옵션 적용 및 증빙 이미지의 엑셀 내 동적 자동 배치 구현.",
        "tasks": ["증빙 이미지 저용량 최적화 프로세스", "XlsxWriter 기반 엑셀 내 위치/크기 자동 지정", "광고주별 증빙 유무/포맷 선택 옵션 개발"],
        "tech": ["Pillow (Image Processing)", "XlsxWriter", "Automated Layout Engine"]
    },
    {
        "title": "[PHASE 3&4] 통합 배포 및 분석 대시보드",
        "date_offset": MILESTONE_DAYS[3],
        "mission": "웹 환경 배포 및 과거 데이터 비교를 통한 수익성 관리 대시보드 구축.",
        "tasks": ["Streamlit 기반 전사 웹 솔루션 배포", "견적 히스토리 추적 및 유사 캠페인 단가 비교", "외주비 세이브율 및 정산 현황 시각화"],
        "tech": ["Streamlit", "SQL/Database", "Data Visualization"]
    }
]

prs = Presentation()

# 1. 표지 슬라이드
slide_layout_title = prs.slide_layouts[0]
slide_title = prs.slides.add_slide(slide_layout_title)
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]
title.text = "AI TF - pioneer"
subtitle.text = "AE 업무 자동화를 위한 견적/증빙 통합 솔루션 로드맵"

title.text_frame.paragraphs[0].font.name = FONT_BOLD
title.text_frame.paragraphs[0].font.size = Pt(44)
subtitle.text_frame.paragraphs[0].font.name = FONT_MEDIUM
subtitle.text_frame.paragraphs[0].font.size = Pt(22)

# 2. 마일스톤 차트 슬라이드
slide_layout_blank = prs.slide_layouts[6]
slide_chart = prs.slides.add_slide(slide_layout_blank)

# 메인 타이틀
title_shape = slide_chart.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
p = title_shape.text_frame.paragraphs[0]
p.text = "PROJECT MILESTONE: AE WORKFLOW AUTOMATION"
p.font.name = FONT_BOLD
p.font.size = Pt(26)
p.font.color.rgb = RGBColor(10, 40, 100)

# 타임라인 화살표
arrow_start_y = Inches(3.2)
shape_arrow = slide_chart.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.2), arrow_start_y, Inches(9.6), Inches(1.2))
shape_arrow.fill.solid()
shape_arrow.fill.fore_color.rgb = RGBColor(230, 240, 250)
shape_arrow.line.color.rgb = RGBColor(100, 150, 200)

colors = [RGBColor(0, 150, 136), RGBColor(3, 169, 244), RGBColor(103, 58, 183), RGBColor(63, 81, 181)]

from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

for i, ms in enumerate(milestones):
    left_x = Inches(0.3 + (i * 2.4))

    # 상단 정보
    txt_box = slide_chart.shapes.add_textbox(left_x, Inches(1.2), Inches(2.2), Inches(1.8))
    tf = txt_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = ms["title"]
    p1.font.name = FONT_BOLD
    p1.font.size = Pt(13)
    p1.font.color.rgb = colors[i]

    target_date = START_DATE + datetime.timedelta(days=ms["date_offset"])
    p2 = tf.add_paragraph()
    p2.text = f"Due: {target_date.strftime('%Y-%m-%d')}"
    p2.font.name = FONT_MEDIUM
    p2.font.size = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = ms["mission"]
    p3.font.name = FONT_LIGHT
    p3.font.size = Pt(9)

    # 하단 세부 과제
    btm_box = slide_chart.shapes.add_textbox(left_x, Inches(4.6), Inches(2.3), Inches(2.5))
    btf = btm_box.text_frame
    btf.word_wrap = True

    p4 = btf.paragraphs[0]
    p4.text = "Key Tasks:"
    p4.font.name = FONT_MEDIUM
    p4.font.size = Pt(10)
    p4.font.bold = True

    for task in ms["tasks"]:
        tp = btf.add_paragraph()
        tp.text = f"• {task}"
        tp.font.name = FONT_LIGHT
        tp.font.size = Pt(8.5)

    # 연결선 (connector 대신 textbox로 대체 - python-pptx 1.x connector API 변경)
    line_left = left_x + Inches(1.1)
    line_box = slide_chart.shapes.add_textbox(line_left, Inches(2.8), Inches(0.01), Inches(0.4))
    line_box.line.color.rgb = colors[i]
    line_box.line.width = Pt(1.5)

output_path = "AI_TF_Pioneer_Roadmap_v2.pptx"
prs.save(output_path)
print(f"파일 생성 완료: {output_path}")
